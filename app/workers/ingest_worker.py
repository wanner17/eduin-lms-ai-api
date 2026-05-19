import os
from app.core.config import settings
from app.core.redis_client import set_ingest_status
from ai_core.rag.chunker import HierarchicalChunker, Page
from ai_core.vectorstore.qdrant_wrapper import QdrantWrapper


def _parse_pdf(file_path: str) -> list[Page]:
    import fitz
    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            pages.append(Page(number=i + 1, text=text))
    return pages


def _parse_pptx(file_path: str) -> list[Page]:
    from pptx import Presentation
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        title = ""
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if shape.name.lower().startswith("title"):
                title = t
            elif t:
                texts.append(t)
        full_text = "\n".join(filter(None, [title] + texts))
        if full_text:
            pages.append(Page(number=i + 1, text=full_text, section_title=title))
    return pages


def _parse_docx(file_path: str) -> list[Page]:
    from docx import Document
    doc = Document(file_path)
    page_size = 20
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages = []
    for i in range(0, len(paragraphs), page_size):
        chunk = paragraphs[i:i + page_size]
        pages.append(Page(number=(i // page_size) + 1, text="\n".join(chunk)))
    return pages


def _parse_file(file_path: str, file_type: str) -> list[Page]:
    parsers = {"pdf": _parse_pdf, "pptx": _parse_pptx, "docx": _parse_docx}
    parser = parsers.get(file_type.lower())
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")
    return parser(file_path)


async def run_ingest(
    material_id: str,
    file_path: str,
    file_type: str,
    file_name: str,
    course_id: int,
    lecture_id: int | None,
):
    await set_ingest_status(material_id, "PROCESSING")
    try:
        pages = _parse_file(file_path, file_type)

        chunker = HierarchicalChunker()
        chunks = chunker.chunk(pages, material_id=material_id)

        from app.core.deps import get_embedding_client
        embed_client = get_embedding_client()
        embeddings = await embed_client.embed([c.text for c in chunks])

        qdrant = QdrantWrapper(
            url=settings.QDRANT_URL,
            collection=settings.QDRANT_COLLECTION,
            vector_size=settings.EMBEDDING_DIM,
        )
        qdrant.ensure_collection()
        qdrant.upsert_chunks(
            chunks=chunks,
            embeddings=embeddings,
            extra_metadata={
                "course_id": course_id,
                "lecture_id": lecture_id,
                "file_name": file_name,
            },
        )

        await set_ingest_status(
            material_id, "READY",
            chunk_count=len(chunks),
            file_name=file_name,
        )

    except Exception as e:
        await set_ingest_status(material_id, "FAILED", error=str(e)[:500])
        raise
    finally:
        # 처리 완료 후 업로드 파일 삭제 (선택)
        if os.path.exists(file_path):
            pass  # 필요 시 os.remove(file_path)
